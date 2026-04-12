"""
PPTX chart data direct-read hook (#1 + #10).

Why this hook exists
--------------------
Docling's PPTX pipeline exports chart visuals as images, which then flow
through the Vision AI for description. Vision can hallucinate numeric
values from blurry chart pixels. When the PPT is authored in PowerPoint,
its charts are already structured data (category / series / values) sitting
in the .pptx zip — reading them directly via python-pptx is 100% accurate,
free, and instant.

What we do
----------
  * Open the original .pptx with python-pptx.
  * Walk every slide in (top, left) reading order (#10 — python-pptx
    defaults to XML-document order which can scramble reading order for
    slides designed by dragging shapes).
  * For every shape where shape.has_chart is True, read the chart's
    categories and series via the python-pptx object model and format
    it as a Markdown table with a "### Chart" heading.
  * Also recurse into GROUP shapes, which can contain charts.
  * Inject the extracted tables into parse_result.markdown at the
    corresponding slide's pagebreak section. The markdown from Docling
    is split by PAGEBREAK_MARKER (one section per slide); we find the
    N-th section (0-indexed = slide N) and append the chart table(s)
    to that section's text.
  * If Docling output has no pagebreak markers (shouldn't happen for PPT
    but we don't crash), fall back to appending all chart tables at the
    document end.
  * Record the number of extracted charts in
    metadata["pptx_charts_extracted"] so the Vision prompt can reference
    it and avoid re-describing chart data it already has in text form.

What we DO NOT do
-----------------
  * Do NOT suppress Vision. Vision still runs and should still describe
    the slide visuals — it just shouldn't re-read chart numbers when
    structured data was already extracted.
  * Do NOT touch slides without charts. Zero cost for non-chart PPTX.
  * Do NOT raise on unsupported chart types (Waterfall / Sunburst /
    Histogram / Treemap use cx:chart namespace and python-pptx raises
    ValueError("unsupported plot type")); these fall through to the
    existing Vision path.

Config:
  parsing.pptx.extract_chart_data — master toggle (default true)
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from ..config import get_nested
from ..parsers.base import ParseResult, PAGEBREAK_MARKER

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Chart formatting
# ---------------------------------------------------------------------------

def _safe_get_axis_title(axis: Any) -> str:
    """
    Extract an axis title if present. python-pptx exposes axis_title on
    CategoryAxis / ValueAxis only when has_title is True, and accessing it
    can raise on some chart types — so we wrap the whole thing.
    """
    try:
        if axis is None or not getattr(axis, "has_title", False):
            return ""
        title = axis.axis_title.text_frame.text
        return title.strip() if title else ""
    except Exception:
        return ""


def _format_chart_as_markdown(chart: Any) -> str | None:
    """
    Convert a python-pptx Chart object into a Markdown block.

    Extracts as much structured context as python-pptx can provide, so the
    downstream Vision step can treat this as ground truth and focus on
    visual elements outside the chart object (annotations, highlights,
    neighbouring text).

    Captures:
      * Chart title
      * Chart type (COLUMN_CLUSTERED, LINE, PIE, …)
      * Value-axis title (e.g. "Revenue (¥M)")
      * Category-axis title (e.g. "Fiscal Quarter")
      * Legend presence
      * Categories, series names, values (the core data table)

    Returns None if the chart type is unsupported (python-pptx raises
    ValueError("unsupported plot type") for Office 2016+ chart types like
    Waterfall / Sunburst / Histogram / Treemap) or if the chart has no
    usable data — both cases fall back to Vision-only handling.

    Format:
        ### Chart: <optional title>
        > type: COLUMN_CLUSTERED · value axis: Revenue (¥M) · legend: yes

        | Category | Series A | Series B |
        |----------|----------|----------|
        | Q1       | 10.5     | 8.2      |
    """
    try:
        plots = chart.plots
        if not plots:
            return None
        plot = plots[0]
        # plot.categories can be empty (e.g. pivot chart without categories)
        categories = list(plot.categories) if plot.categories is not None else []
        series_list = list(chart.series)
    except ValueError as e:
        # "unsupported plot type" — Office 2016 chart types
        logger.debug(f"Chart skipped (unsupported plot type): {e}")
        return None
    except Exception as e:
        logger.debug(f"Chart read failed: {e}")
        return None

    if not categories or not series_list:
        return None

    # --- Optional metadata (all best-effort, any failure leaves field empty) ---
    title = ""
    try:
        if chart.has_title and chart.chart_title.has_text_frame:
            title = chart.chart_title.text_frame.text.strip()
    except Exception:
        pass

    chart_type = ""
    try:
        ct = chart.chart_type
        # Enum values expose .name (e.g. "COLUMN_CLUSTERED")
        chart_type = getattr(ct, "name", str(ct)) if ct is not None else ""
    except Exception:
        pass

    value_axis_title = _safe_get_axis_title(getattr(chart, "value_axis", None))
    category_axis_title = _safe_get_axis_title(getattr(chart, "category_axis", None))

    has_legend = False
    try:
        has_legend = bool(chart.has_legend)
    except Exception:
        pass

    # Build context line only if at least one field is populated
    context_bits: list[str] = []
    if chart_type:
        context_bits.append(f"type: {chart_type}")
    if value_axis_title:
        context_bits.append(f"value axis: {value_axis_title}")
    if category_axis_title:
        context_bits.append(f"category axis: {category_axis_title}")
    if has_legend:
        context_bits.append("legend: yes")

    # --- Build 2D data table ---
    header = ["Category"] + [_series_name(s, i) for i, s in enumerate(series_list)]
    data_rows: list[list[str]] = []

    for idx, cat in enumerate(categories):
        row = [_cell_str(cat)]
        for series in series_list:
            try:
                values = list(series.values) if series.values else []
                value = values[idx] if idx < len(values) else ""
            except Exception:
                value = ""
            row.append(_cell_str(value))
        data_rows.append(row)

    # --- Render Markdown ---
    lines: list[str] = []
    lines.append(f"### Chart{': ' + title if title else ''}")
    if context_bits:
        # Use a blockquote so it renders nicely but is visually distinct
        # from the data table
        lines.append(f"> {' · '.join(context_bits)}")
    lines.append("")
    lines.append("| " + " | ".join(header) + " |")
    lines.append("|" + "|".join(["---"] * len(header)) + "|")
    for row in data_rows:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _series_name(series: Any, default_idx: int) -> str:
    """Get a series display name, falling back to 'Series N' if absent."""
    try:
        name = series.name
        if name:
            return str(name).strip()
    except Exception:
        pass
    return f"Series {default_idx + 1}"


def _cell_str(value: Any) -> str:
    """
    Convert a cell value to a display string.

    - None → ""
    - floats with no fractional part → int-like display (1.0 → "1")
    - other numbers → str(value)
    - strings → stripped
    - escapes pipe characters to avoid breaking the table
    """
    if value is None:
        return ""
    if isinstance(value, float):
        if value.is_integer():
            return str(int(value))
        return f"{value:g}"
    return str(value).replace("|", "\\|").strip()


# ---------------------------------------------------------------------------
# Slide traversal (reading-order-aware, recurses into group shapes)
# ---------------------------------------------------------------------------

def _iter_shapes_in_reading_order(shapes: Any) -> list[Any]:
    """
    Sort shapes by (top, left) to approximate reading order (#10).

    python-pptx returns shapes in XML document order, which is usually the
    order they were added to the slide — not the visual reading order.
    Sorting by top-then-left gives a reasonable left-to-right top-to-bottom
    traversal. Shapes with unset top/left (rare, e.g. placeholders) sort
    to the start.
    """
    def sort_key(shape: Any) -> tuple[int, int]:
        top = shape.top if shape.top is not None else 0
        left = shape.left if shape.left is not None else 0
        return (int(top), int(left))

    return sorted(shapes, key=sort_key)


def _collect_charts_on_slide(slide: Any) -> list[str]:
    """
    Walk a slide's shape tree in reading order and collect rendered chart
    Markdown tables. Recurses into GROUP shapes. Returns a list of
    Markdown strings (one per chart).
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-not-found]

    charts_md: list[str] = []

    def walk(shapes: Any) -> None:
        for shape in _iter_shapes_in_reading_order(shapes):
            if getattr(shape, "has_chart", False):
                md = _format_chart_as_markdown(shape.chart)
                if md:
                    charts_md.append(md)
            # Recurse into groups
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    walk(shape.shapes)
                except Exception as e:
                    logger.debug(f"Failed to recurse group shape: {e}")

    walk(slide.shapes)
    return charts_md


# ---------------------------------------------------------------------------
# Markdown injection
# ---------------------------------------------------------------------------

def _inject_into_markdown(
    parse_result: ParseResult,
    per_slide_charts: list[list[str]],
) -> int:
    """
    Insert chart Markdown blocks into the corresponding pagebreak sections.

    Returns the total number of charts actually injected. Charts that
    cannot be aligned (slide index beyond section count) fall through to
    appended-at-end behavior.
    """
    if not any(per_slide_charts):
        return 0

    sections = parse_result.markdown.split(PAGEBREAK_MARKER)
    has_pagebreaks = len(sections) > 1
    injected = 0

    if has_pagebreaks:
        # Align slide index → section index. Docling's PPT pipeline produces
        # one section per slide, so slide N (0-indexed) maps to sections[N].
        for slide_idx, charts in enumerate(per_slide_charts):
            if not charts:
                continue
            if slide_idx < len(sections):
                chart_block = "\n\n" + "\n\n".join(charts)
                sections[slide_idx] = sections[slide_idx].rstrip() + chart_block + "\n"
                injected += len(charts)
        parse_result.markdown = PAGEBREAK_MARKER.join(sections)
    else:
        # Fallback: no pagebreaks in Docling output. Append all charts at end.
        flat = [c for charts in per_slide_charts for c in charts]
        if flat:
            parse_result.markdown = (
                parse_result.markdown.rstrip()
                + "\n\n"
                + "\n\n".join(flat)
                + "\n"
            )
            injected = len(flat)

    return injected


# ---------------------------------------------------------------------------
# Hook entry point
# ---------------------------------------------------------------------------

def pptx_chart_hook(
    file_path: Path,
    parse_result: ParseResult,
    config: dict[str, Any],
) -> None:
    """
    Extract structured chart data from a .pptx file and inject into markdown.

    Runs in the "post_parse" phase (after Docling, before Vision). Never
    raises — any failure leaves parse_result untouched and logs a warning.
    """
    if not get_nested(config, "parsing.pptx.extract_chart_data", True):
        return

    # Only pptx (not legacy .ppt — python-pptx cannot read binary .ppt)
    if file_path.suffix.lower() != ".pptx":
        return

    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError:
        logger.debug(
            "python-pptx not installed; skipping PPTX chart extraction. "
            "pip install python-pptx to enable."
        )
        return

    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        logger.warning(f"python-pptx failed to open {file_path.name}: {e}")
        return

    per_slide_charts: list[list[str]] = []
    for slide in prs.slides:
        try:
            charts = _collect_charts_on_slide(slide)
        except Exception as e:
            logger.debug(f"Chart extraction failed on slide: {e}")
            charts = []
        per_slide_charts.append(charts)

    injected = _inject_into_markdown(parse_result, per_slide_charts)

    # Publish per-page structured extractions so the Vision step can
    # reference them as ground truth and focus on visual context Vision
    # can see but python-pptx can't (annotations, highlights, callouts).
    # Key: 1-based page number (matches PageData.page_no). Value: the
    # full Markdown block(s) already injected into parse_result.markdown.
    #
    # This dict is shared across hooks — other extractors (future) can
    # merge their own entries by extending per-page strings.
    if injected > 0:
        extractions: dict[int, str] = parse_result.metadata.setdefault(
            "structured_extractions_per_page", {}
        )
        for slide_idx, charts in enumerate(per_slide_charts):
            if not charts:
                continue
            page_no = slide_idx + 1  # PageData.page_no is 1-based
            merged = "\n\n".join(charts)
            if page_no in extractions:
                extractions[page_no] = extractions[page_no] + "\n\n" + merged
            else:
                extractions[page_no] = merged

        parse_result.metadata["pptx_charts_extracted"] = injected
        logger.info(
            f"PPTX chart extraction: {injected} chart(s) injected from {file_path.name}"
        )
