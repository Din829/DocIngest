"""
Document inspection — lightweight pre-flight check before pipeline processing.

Quickly scans input files to report size, page count, format, and estimated
processing characteristics WITHOUT actually parsing them. Designed to be
called by Agents, MCP tools, or humans before ``docingest run``.

Usage:
    CLI:     docingest inspect report.pdf slides.pptx ./docs/
    Python:  from docingest.inspect import inspect_files
             results = inspect_files([Path("report.pdf")])

Performance: sub-second for most files (reads only metadata, not content).
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from .config import get_nested

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Per-format inspectors (fast metadata only, no full parsing)
# ---------------------------------------------------------------------------
# All inspectors take (file_path, config) and return a dict. Keys commonly
# produced:
#   pages          — page/slide/sheet count (whatever "unit of processing" is)
#   chars_est      — estimated text length (proxy for output bulk)
#   total_rows     — xlsx-specific
#   duration_sec   — audio/video duration from ffprobe
#   words          — docx-specific
#   error          — exception message if introspection failed
#
# A returned {} or a dict missing a field simply means "no signal for that
# dimension" — the safety layer treats absent metrics as "don't check".


def _pick_sample_indices(total: int, n: int) -> list[int]:
    """
    Pick up to n representative page indices for sampling.
    Always includes first / last when total >= 2, fills middle evenly.
    """
    if total <= 0 or n <= 0:
        return []
    if total <= n:
        return list(range(total))
    if n == 1:
        return [total // 2]
    # First, evenly spaced middle picks, last
    picks = [0, total - 1]
    if n >= 3:
        mid_slots = n - 2
        for k in range(1, mid_slots + 1):
            picks.append(k * total // (mid_slots + 1))
    return sorted(set(picks))


def _inspect_pdf(file_path: Path, config: dict[str, Any]) -> dict[str, Any]:
    """
    PDF: page count (O(1) via PyMuPDF xref) + sampled text-length estimate.

    Sampling keeps inspect fast on large PDFs — full get_text() on a
    1000-page document costs ~10 s; three-page sampling is ~100 ms.
    """
    try:
        import pymupdf
        doc = pymupdf.open(str(file_path))
        pages = len(doc)
        chars_est = 0
        if pages > 0:
            sampling_enabled = get_nested(config, "safety.sampling.enabled", True)
            if sampling_enabled:
                n = int(get_nested(config, "safety.sampling.pdf_sample_pages", 3))
                indices = _pick_sample_indices(pages, n)
                if indices:
                    sample_total = sum(
                        len(doc.load_page(i).get_text() or "") for i in indices
                    )
                    avg = sample_total / len(indices)
                    chars_est = int(avg * pages)
            else:
                # Exact — slower but precise.
                chars_est = sum(
                    len(doc.load_page(i).get_text() or "") for i in range(pages)
                )
        doc.close()
        return {"pages": pages, "chars_est": chars_est}
    except Exception as e:
        logger.debug(f"PDF inspection failed for {file_path.name}: {e}")
        return {"pages": None, "error": str(e)}


def _inspect_pptx(file_path: Path, config: dict[str, Any]) -> dict[str, Any]:
    """PPTX: slide count + aggregated text from every shape.text_frame."""
    _ = config  # sampling not applicable — slide counts are small
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        slides = list(prs.slides)
        chars_est = 0
        for slide in slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    try:
                        chars_est += len(shape.text_frame.text or "")
                    except Exception:
                        continue
        return {"pages": len(slides), "chars_est": chars_est}
    except Exception as e:
        logger.debug(f"PPTX inspection failed for {file_path.name}: {e}")
        return {"pages": None, "error": str(e)}


def _inspect_xlsx(file_path: Path, config: dict[str, Any]) -> dict[str, Any]:
    """XLSX: sheet count + total row count (char estimation skipped —
    total_rows is the stronger processing-cost signal for spreadsheets)."""
    _ = config
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        sheets = wb.sheetnames
        total_rows = 0
        for name in sheets:
            ws = wb[name]
            total_rows += ws.max_row or 0
        wb.close()
        return {"pages": len(sheets), "sheets": sheets, "total_rows": total_rows}
    except Exception as e:
        logger.debug(f"XLSX inspection failed for {file_path.name}: {e}")
        return {"pages": None, "error": str(e)}


def _inspect_docx(file_path: Path, config: dict[str, Any]) -> dict[str, Any]:
    """DOCX: approximate page count + character count from paragraphs."""
    _ = config
    try:
        from docx import Document
        doc = Document(str(file_path))
        word_count = 0
        chars_est = 0
        for p in doc.paragraphs:
            text = p.text or ""
            chars_est += len(text)
            word_count += len(text.split())
        est_pages = max(1, word_count // 300)
        return {
            "pages": est_pages,
            "chars_est": chars_est,
            "words": word_count,
            "pages_estimated": True,
        }
    except Exception as e:
        logger.debug(f"DOCX inspection failed for {file_path.name}: {e}")
        return {"pages": None, "error": str(e)}


def _inspect_zip(file_path: Path, config: dict[str, Any]) -> dict[str, Any]:
    """ZIP: entry count + uncompressed size. No page concept."""
    _ = config
    try:
        import zipfile
        with zipfile.ZipFile(str(file_path), "r") as zf:
            entries = zf.infolist()
            file_count = sum(1 for e in entries if not e.is_dir())
            total_size = sum(e.file_size for e in entries)
            return {
                "pages": file_count,
                "files_inside": file_count,
                "uncompressed_size_mb": round(total_size / (1024 * 1024), 1),
            }
    except Exception as e:
        logger.debug(f"ZIP inspection failed for {file_path.name}: {e}")
        return {"pages": None, "error": str(e)}


def _parse_duration_string(s: str) -> int | None:
    """Convert ffprobe's H:MM:SS or MM:SS string to integer seconds."""
    if not s:
        return None
    try:
        parts = s.split(":")
        nums = [int(p) for p in parts]
        if len(nums) == 3:
            return nums[0] * 3600 + nums[1] * 60 + nums[2]
        if len(nums) == 2:
            return nums[0] * 60 + nums[1]
        if len(nums) == 1:
            return nums[0]
    except (ValueError, AttributeError):
        return None
    return None


def _inspect_media(file_path: Path, config: dict[str, Any]) -> dict[str, Any]:
    """
    Audio / video: duration in seconds via ffprobe.

    Reuses MediaParser's _get_duration implementation to avoid duplicating
    the ffprobe invocation. Returns an empty dict (no duration_sec field)
    when ffprobe is unavailable — safety then simply won't flag duration.
    """
    try:
        from .parsers.media_parser import MediaParser
        mp = MediaParser(config)
        duration_str = mp._get_duration(file_path)
        duration_sec = _parse_duration_string(duration_str)
        if duration_sec is None:
            return {}
        return {"duration_sec": duration_sec}
    except Exception as e:
        logger.debug(f"Media inspection failed for {file_path.name}: {e}")
        return {}


def _inspect_text(file_path: Path, config: dict[str, Any]) -> dict[str, Any]:
    """TXT / MD / CSV: byte count as a char_est proxy (UTF-8 ≈ chars)."""
    _ = config
    try:
        size = file_path.stat().st_size
        return {"chars_est": size}
    except OSError:
        return {}


# Format → inspector mapping. Unknown suffixes fall through to size-only.
_INSPECTORS: dict[str, Any] = {
    "pdf": _inspect_pdf,
    "pptx": _inspect_pptx,
    "ppt": _inspect_pptx,
    "xlsx": _inspect_xlsx,
    "xls": _inspect_xlsx,
    "docx": _inspect_docx,
    "doc": _inspect_docx,
    "zip": _inspect_zip,
    # Text-based (cheap byte-count estimate)
    "txt": _inspect_text,
    "md": _inspect_text,
    "csv": _inspect_text,
    "tsv": _inspect_text,
    "log": _inspect_text,
    # Media — duration only, via ffprobe
    "mp3": _inspect_media, "wav": _inspect_media, "m4a": _inspect_media,
    "flac": _inspect_media, "aac": _inspect_media, "ogg": _inspect_media,
    "wma": _inspect_media, "opus": _inspect_media,
    "mp4": _inspect_media, "avi": _inspect_media, "mkv": _inspect_media,
    "webm": _inspect_media, "mov": _inspect_media, "wmv": _inspect_media,
    "flv": _inspect_media, "ts": _inspect_media, "m4v": _inspect_media,
}


# ---------------------------------------------------------------------------
# Recommendation logic (driven by safety.per_file thresholds)
# ---------------------------------------------------------------------------

def _recommend(info: dict[str, Any], config: dict[str, Any]) -> str:
    """
    Generate a human/agent-readable recommendation based on inspection.

    Thresholds come from the safety.* config section — the same values
    Phase 0 uses — so `inspect` and `run` agree on what "too big" means.
    Adding or changing a threshold only touches config.
    """
    # Lazy import to avoid circular dependency; safety.py imports us too.
    from .safety import check_file_violations

    warnings: list[str] = []

    reasons = check_file_violations(info, config)
    for r in reasons:
        val = r.get("value")
        thr = r.get("threshold")
        if isinstance(val, (int, float)) and isinstance(thr, (int, float)):
            warnings.append(f"{r['metric']}={val:,} > {thr:,}")
        else:
            warnings.append(f"{r['metric']}={val} > {thr}")

    if info.get("error"):
        warnings.append(f"Inspection error: {info['error']}")

    if not warnings:
        return "Ready"
    return "; ".join(warnings)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def inspect_single(file_path: Path, config: dict[str, Any] | None = None) -> dict[str, Any]:
    """
    Inspect a single file without parsing it.

    Returns a dict with: name, format, size_mb, pages, est_cost_usd,
    recommendation, and any format-specific fields (sheets, words, chars_est,
    duration_sec, files_inside, etc.). Safe on unknown formats — those get
    size-only info.
    """
    if config is None:
        config = {}

    suffix = file_path.suffix.lstrip(".").lower()
    try:
        size_mb = file_path.stat().st_size / (1024 * 1024)
    except OSError as e:
        size_mb = 0.0
        logger.debug(f"stat() failed for {file_path.name}: {e}")

    info: dict[str, Any] = {
        "name": file_path.name,
        "path": str(file_path),
        "format": suffix,
        "size_mb": round(size_mb, 2),
    }

    # Run format-specific inspector — all inspectors take (path, config)
    # and are free to return any subset of the documented fields.
    inspector = _INSPECTORS.get(suffix)
    if inspector:
        try:
            info.update(inspector(file_path, config))
        except Exception as e:
            # Never fail inspect_single due to a sub-inspector crash.
            logger.debug(f"Inspector for .{suffix} crashed on {file_path.name}: {e}")
            info.setdefault("error", str(e))
    else:
        info["pages"] = None

    # Cost estimation: cheap lookup based on page count + Vision model price.
    # Lazy import so inspect.py stays self-contained for callers that only
    # want size/pages without paying the safety import.
    from .safety import estimate_file_cost_usd
    info["est_cost_usd"] = round(estimate_file_cost_usd(info, config), 4)

    info["recommendation"] = _recommend(info, config)
    return info


def inspect_files(
    input_paths: list[Path | str],
    config: dict[str, Any] | None = None,
) -> list[dict[str, Any]]:
    """
    Inspect multiple files/directories.

    Expands directories recursively (same logic as pipeline.discover_files).
    Returns list of inspection results, one per file.

    Invalid inputs (missing paths, failed URL resolution) are returned as
    inspection entries with ``format="invalid"`` and ``error_reason``/
    ``error_detail`` fields — same channel as valid inspections so callers
    can iterate uniformly. Catching these before `ingest()` is exactly what
    inspect is for: a typo / cross-container mistake / broken URL should
    surface here, not as a silent zero-result run later.
    """
    from .pipeline import discover_files

    if config is None:
        from .config import load_config
        config = load_config()

    # discover_files returns (valid, invalid). Both feed into the returned
    # list so inspect output covers everything the caller asked about.
    files, invalid_inputs = discover_files(
        [Path(p) if isinstance(p, str) else p for p in input_paths],
        config,
    )

    results: list[dict[str, Any]] = [inspect_single(f, config) for f in files]
    for inv in invalid_inputs:
        results.append({
            "name": Path(inv.input).name if inv.input else inv.input,
            "path": inv.input,
            "format": "invalid",
            "error_reason": inv.reason,
            "error_detail": inv.detail,
            "recommendation": (
                f"Input cannot be processed ({inv.reason}). "
                f"Fix this before calling ingest() — otherwise it surfaces "
                f"as a silent failed=1 entry in result.stats['errors']."
            ),
        })
    return results
