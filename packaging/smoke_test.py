"""Host-side smoke driver — the build gate of build_exe.ps1 stage 4.

Generates a tiny deterministic sample of every offline-parseable format,
then runs `DocIngest.exe --smoke <samples> <out>` so the dependency matrix
executes INSIDE the frozen bundle (see launch_gui.py for why that matters).
Exit code is the gate: 0 = ship-able, non-zero = incomplete bundle.

Samples are generated programmatically (not pulled from test_docs/) so the
gate is self-contained, fast, and runs the same on any build machine. The
generator libs (pymupdf / python-docx / python-pptx / openpyxl / PIL) are
exactly the dev deps this repo already requires.

Usage:
    python packaging/smoke_test.py <dist_dir>     # e.g. packaging/dist/DocIngest
"""
from __future__ import annotations

import subprocess
import sys
import tempfile
import time
import zipfile
from pathlib import Path


def make_samples(target: Path) -> None:
    """One small real file per offline-parseable format."""
    target.mkdir(parents=True, exist_ok=True)

    # pdf — pymupdf (exercises docling PDF pipeline: layout + tableformer)
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 100), "Smoke Test PDF", fontsize=20)
    page.insert_text((72, 140), "Section content for chunking. 日本語テキスト。")
    doc.save(str(target / "sample.pdf"))
    doc.close()

    # docx — python-docx (docling word backend)
    import docx
    d = docx.Document()
    d.add_heading("Smoke DOCX", level=1)
    d.add_paragraph("Body paragraph with enough text to chunk. 仕様書サンプル。")
    d.save(str(target / "sample.docx"))

    # pptx — python-pptx (docling + LibreOffice page-image path)
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Smoke PPTX"
    slide.placeholders[1].text = "Bullet one\nBullet two"
    prs.save(str(target / "sample.pptx"))

    # xlsx — openpyxl (the openpyxl renderer path)
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "data"
    ws.append(["id", "name", "値"])
    ws.append([1, "alpha", 10])
    ws.append([2, "beta", 20])
    wb.save(str(target / "sample.xlsx"))

    # png — PIL (image input route)
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (400, 200), "white")
    ImageDraw.Draw(img).text((20, 80), "Smoke PNG", fill="black")
    img.save(str(target / "sample.png"))

    # md — text route
    (target / "sample.md").write_text(
        "# Smoke MD\n\nPlain markdown passes through the text parser.\n",
        encoding="utf-8",
    )

    # zip — archive expansion route (wraps the md)
    with zipfile.ZipFile(target / "sample.zip", "w") as z:
        z.write(target / "sample.md", "inner/zipped.md")


def main() -> int:
    if len(sys.argv) != 2:
        print(__doc__)
        return 2
    dist = Path(sys.argv[1]).resolve()
    exe = dist / "DocIngest.exe"
    if not exe.is_file():
        print(f"error: {exe} not found — build first.")
        return 1

    with tempfile.TemporaryDirectory(prefix="docingest_smoke_") as td:
        samples = Path(td) / "samples"
        out = Path(td) / "kb"
        print(f"Generating samples → {samples}")
        make_samples(samples)

        print(f"Running in-bundle matrix: {exe} --smoke")
        t0 = time.perf_counter()
        proc = subprocess.run(
            [str(exe), "--smoke", str(samples), str(out)],
            timeout=1800,
        )
        elapsed = time.perf_counter() - t0

        # The exe is windowed (no stdout) — the matrix writes its report to
        # <out>/smoke.log instead; relay it here for the build console / CI.
        log = out / "smoke.log"
        if log.is_file():
            print("\n----- in-bundle smoke.log -----")
            print(log.read_text(encoding="utf-8"))
            print("----- end smoke.log -----")
        else:
            print("warning: no smoke.log produced — exe likely crashed before the matrix started.")

        print(f"\nIn-bundle matrix finished in {elapsed:.1f}s, exit={proc.returncode}")
        return proc.returncode


if __name__ == "__main__":
    sys.exit(main())
